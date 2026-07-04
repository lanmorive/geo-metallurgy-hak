"""Парсинг PPTX/POTX через python-pptx."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.shapes.base import BaseShape

from app.ingest.authors import fix_file_metadata_author
from app.ingest.table_md import dedup_row_cells, rows_to_markdown
from app.ingest.types import Block, DocMeta, ParseResult

_TITLE_PLACEHOLDERS = frozenset({PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE})


def _iter_shapes(container) -> list[BaseShape]:
  """Рекурсивно собрать shapes (включая вложенные в GroupShape)."""
  shapes: list[BaseShape] = []
  for shape in container.shapes:
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
      shapes.extend(_iter_shapes(shape))
    else:
      shapes.append(shape)
  return shapes


def _sort_shapes(shapes: list[BaseShape]) -> list[BaseShape]:
  return sorted(shapes, key=lambda s: (s.top, s.left))


def _pptx_table_to_markdown(table) -> str:
  parsed_rows: list[list[str]] = []
  for row in table.rows:
    cells = dedup_row_cells([cell.text.strip() for cell in row.cells])
    parsed_rows.append(cells)
  return rows_to_markdown(parsed_rows)


def _shape_to_blocks(shape: BaseShape, slide_num: int) -> list[Block]:
  blocks: list[Block] = []

  if shape.has_table:
    md = _pptx_table_to_markdown(shape.table)
    if md:
      blocks.append(Block(type="table", text=md, page=slide_num))
    return blocks

  if not shape.has_text_frame:
    return blocks

  text = shape.text_frame.text.strip()
  if not text:
    return blocks

  if shape.is_placeholder:
    ph_type = shape.placeholder_format.type
    if ph_type in _TITLE_PLACEHOLDERS:
      blocks.append(
        Block(type="heading", text=text, level=2, page=slide_num)
      )
      return blocks

  blocks.append(Block(type="paragraph", text=text, page=slide_num))
  return blocks


def _notes_block(slide, slide_num: int) -> Block | None:
  try:
    notes_slide = slide.notes_slide
  except ValueError:
    return None
  if notes_slide is None:
    return None
  notes_frame = notes_slide.notes_text_frame
  if notes_frame is None:
    return None
  text = notes_frame.text.strip()
  if not text:
    return None
  return Block(
    type="paragraph",
    text=f"[заметки] {text}",
    page=slide_num,
  )


def parse_pptx(path: Path) -> ParseResult:
  presentation = Presentation(str(path))
  blocks: list[Block] = []
  markdown_parts: list[str] = []
  image_only_slides = 0

  props = presentation.core_properties
  created = props.created.isoformat() if props.created else None
  title = props.title.strip() if props.title and props.title.strip() else None
  doc_meta = DocMeta(
    file_metadata_author=fix_file_metadata_author(props.author),
    file_metadata_title=title,
    created=created,
    pages=len(presentation.slides),
    image_only_slides=0,
  )

  for slide_num, slide in enumerate(presentation.slides, start=1):
    slide_blocks: list[Block] = []

    for shape in _sort_shapes(_iter_shapes(slide)):
      slide_blocks.extend(_shape_to_blocks(shape, slide_num))

    notes = _notes_block(slide, slide_num)
    if notes is not None:
      slide_blocks.append(notes)

    if not slide_blocks:
      image_only_slides += 1
      continue

    for block in slide_blocks:
      blocks.append(block)
      if block.type == "heading":
        markdown_parts.append(f"{'#' * (block.level or 2)} {block.text}")
      else:
        markdown_parts.append(block.text)

  doc_meta.image_only_slides = image_only_slides

  return ParseResult(
    markdown_text="\n\n".join(markdown_parts),
    blocks=blocks,
    doc_meta=doc_meta,
  )
